#!/usr/bin/env python3
"""Build the Melo investor pitch deck in PowerPoint, branded to match the app."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand palette ────────────────────────────────────────────────────────────
MOSS         = RGBColor(0x26, 0x42, 0x36)  # primary dark green
MOSS_LIGHT   = RGBColor(0x3F, 0x6B, 0x4E)
SAGE         = RGBColor(0x7A, 0x9C, 0x82)
SAGE_LIGHT   = RGBColor(0xE4, 0xEB, 0xE3)
CLAY         = RGBColor(0xC5, 0x7B, 0x57)  # warm terracotta accent
CLAY_LIGHT   = RGBColor(0xF1, 0xE0, 0xD6)
CREAM        = RGBColor(0xF6, 0xF1, 0xE7)  # body background
WARM_WHITE   = RGBColor(0xFB, 0xF8, 0xF2)
BARK         = RGBColor(0x4E, 0x3C, 0x2C)  # dark text on light bg
STONE        = RGBColor(0x9A, 0x95, 0x8B)
TEXT_LIGHT   = RGBColor(0xE6, 0xDC, 0xC5)  # cream text on dark bg

SERIF = "Georgia"
SANS  = "Calibri"

# Brand assets — cleaned versions of the Melo app icon
LOGO_BADGE = "/sessions/gracious-lucid-hopper/mnt/outputs/melo-pitch/melo-logo-badge.png"
LOGO_HERO  = "/sessions/gracious-lucid-hopper/mnt/outputs/melo-pitch/melo-logo-hero.png"

# ── Setup ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title  = "Melo — Investor Pitch"
prs.core_properties.author = "Melo Parents ApS"

BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_bg(slide, color):
    """Full-bleed background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_text(slide, text, x, y, w, h, *,
             font=SANS, size=14, color=BARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a single-paragraph text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb

def add_lines(slide, lines, x, y, w, h, *,
              font=SANS, size=14, color=BARK, bullet=False,
              align=PP_ALIGN.LEFT, line_spacing=1.25, gap_after=4):
    """Add a multi-line text box. `lines` = list of strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(gap_after)
        run = p.add_run()
        run.text = ("• " + line) if bullet else line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def add_card(slide, x, y, w, h, fill=WARM_WHITE, line=None):
    """Rounded card with optional border."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.08
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s

def add_accent_dot(slide, x, y, size=Inches(0.15), color=CLAY):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def add_page_number(slide, n, total):
    add_text(slide, f"{n} / {total}", Inches(12.6), Inches(7.15), Inches(0.6), Inches(0.3),
             size=9, color=STONE, align=PP_ALIGN.RIGHT)

def add_brand_mark(slide, on_dark=False):
    """Small icon + wordmark in the top-left corner of content slides."""
    color = TEXT_LIGHT if on_dark else MOSS
    # Logo badge — clean cropped icon
    slide.shapes.add_picture(LOGO_BADGE, Inches(0.55), Inches(0.3),
                             height=Inches(0.55), width=Inches(0.55))
    add_text(slide, "Melo", Inches(1.25), Inches(0.32), Inches(2), Inches(0.4),
             font=SERIF, size=20, color=color, bold=True)
    add_text(slide, "FOR NEW PARENTS",
             Inches(1.25), Inches(0.65), Inches(3), Inches(0.25),
             size=8, color=color)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, MOSS)

# Decorative accent dots scattered subtly
for (x, y) in [(11.2, 1.0), (11.8, 1.4), (12.4, 1.1), (0.6, 6.4), (1.2, 6.7)]:
    add_accent_dot(s, Inches(x), Inches(y), size=Inches(0.12), color=CLAY)

# Hero logo top-right
s.shapes.add_picture(LOGO_HERO, Inches(9.4), Inches(0.65),
                     height=Inches(2.4), width=Inches(2.4))

# Brand wordmark top-left
add_text(s, "Melo", Inches(0.8), Inches(0.7), Inches(3), Inches(0.7),
         font=SERIF, size=36, color=TEXT_LIGHT, bold=True)
add_text(s, "FOR NEW PARENTS", Inches(0.8), Inches(1.25), Inches(4), Inches(0.4),
         size=10, color=TEXT_LIGHT)

add_text(s, "Pregnancy. Parenthood.\nFor two parents.",
         Inches(0.8), Inches(3.0), Inches(11), Inches(2.6),
         font=SERIF, size=60, color=TEXT_LIGHT, bold=False, line_spacing=1.05)

add_text(s, "Investor pitch — Marbella 2026",
         Inches(0.8), Inches(5.6), Inches(8), Inches(0.4),
         size=14, color=TEXT_LIGHT)
add_text(s, "Louise Bencard, Founder · Melo Parents ApS",
         Inches(0.8), Inches(6.0), Inches(8), Inches(0.4),
         size=12, color=TEXT_LIGHT, italic=True)
add_text(s, "meloparents.com",
         Inches(0.8), Inches(6.4), Inches(8), Inches(0.4),
         size=11, color=TEXT_LIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Hook
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

# Large stat
add_text(s, "67%",
         Inches(0.8), Inches(2.0), Inches(7), Inches(2.5),
         font=SERIF, size=180, color=MOSS, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, "of couples experience a significant decline in",
         Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5),
         font=SERIF, size=24, color=BARK)
add_text(s, "relationship satisfaction in the first year of parenthood.",
         Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
         font=SERIF, size=24, color=BARK)

add_text(s, "Equitable task-sharing is the strongest predictor of recovery.",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
         size=14, color=CLAY, italic=True)
add_text(s, "Source: Gottman Institute, longitudinal couples research",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.3),
         size=10, color=STONE, italic=True)

add_page_number(s, 2, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Problem
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "The problem",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Three gaps in how new families are served today",
         Inches(0.8), Inches(1.85), Inches(10), Inches(0.4),
         size=14, color=STONE, italic=True)

problems = [
    ("01", "The father gap",
     "Pregnancy and parenting content is written for mothers. Fathers install the app, find nothing for them, and stop opening it within four weeks."),
    ("02", "The mental load gap",
     "Even when fathers are engaged, the invisible work of running a household falls on the mother. A leading driver of postpartum depression and relationship breakdown."),
    ("03", "The trust gap",
     "The market is dominated by ad-supported and data-monetized models. Vulnerable, sleep-deprived parents are sold supplements and lifestyle products at scale."),
]
for i, (num, title, body) in enumerate(problems):
    y = Inches(2.6 + i * 1.45)
    add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.25), fill=WARM_WHITE, line=SAGE_LIGHT)
    add_text(s, num, Inches(1.0), y + Inches(0.18), Inches(0.9), Inches(0.5),
             font=SERIF, size=28, color=CLAY, bold=True)
    add_text(s, title, Inches(1.9), y + Inches(0.2), Inches(4), Inches(0.5),
             font=SERIF, size=20, color=MOSS, bold=True)
    add_text(s, body, Inches(1.9), y + Inches(0.65), Inches(10.4), Inches(0.55),
             size=12.5, color=BARK, line_spacing=1.25)

add_page_number(s, 3, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — The Solution
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "The solution",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "One app, two parents, role-aware throughout.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# Two-column: mother | partner — joined by a center 'AND' column
col_w = Inches(5.4)
col_h = Inches(4.2)
y0    = Inches(2.6)

# Mother column
add_card(s, Inches(0.8), y0, col_w, col_h, fill=CLAY_LIGHT)
add_text(s, "For the mother",
         Inches(1.0), y0 + Inches(0.3), col_w - Inches(0.4), Inches(0.5),
         font=SERIF, size=20, color=BARK, bold=True)
add_lines(s, [
    "Week-by-week pregnancy verified against Danish health authority",
    "Postpartum recovery support",
    "Sleep, diary, and milestone tracking",
    "Personalized nutrition and supplement guidance",
    "Reminders that she also needs care",
], Inches(1.0), y0 + Inches(0.95), col_w - Inches(0.4), Inches(3.0),
   size=12, color=BARK, line_spacing=1.4, gap_after=4)

# Partner column
add_card(s, Inches(7.1), y0, col_w, col_h, fill=SAGE_LIGHT)
add_text(s, "For the father / co-parent",
         Inches(7.3), y0 + Inches(0.3), col_w - Inches(0.4), Inches(0.5),
         font=SERIF, size=20, color=MOSS, bold=True)
add_lines(s, [
    "Daily missions and content in a language that fits him",
    "Concrete ways to support mom, adapted to feeding method",
    "Age-aware play ideas (yes, play = workout)",
    "Kettlebell-weight comparisons, achievement framing",
    "Built-in reminders to verbalize appreciation",
], Inches(7.3), y0 + Inches(0.95), col_w - Inches(0.4), Inches(3.0),
   size=12, color=BARK, line_spacing=1.4, gap_after=4)

# Connection layer
add_text(s, "Connected by",
         Inches(0.8), Inches(7.05), Inches(2), Inches(0.3),
         size=10, color=STONE, italic=True)
add_text(s, "Take-a-task   ·   Send appreciation   ·   Shared needs   ·   Date-night dreams   ·   MELO AI",
         Inches(0.8), Inches(7.3), Inches(11.5), Inches(0.3),
         size=11, color=MOSS, bold=True)

add_page_number(s, 4, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Product
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Product",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Built and ready. Six features that define the experience.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

features = [
    ("📅", "Week-by-week", "From the first heartbeat to term. Verified content, no surprises."),
    ("😴", "Sleep tracker", "One-tap stopwatch. See rhythms over weeks. Find calm on hard nights."),
    ("📔", "Shared diary", "Feeds, diapers, sleep, milestones — both parents log and sync."),
    ("🤝", "Take-and-thank", "Take a task off your partner. Send a one-tap appreciation back."),
    ("💬", "MELO AI", "Ask anything. Grounded in Danish professional guidance. Refers when needed."),
    ("💪", "Father-aware", "Dedicated dashboard, age-aware content, support-mom missions."),
]
gx0 = Inches(0.8); gy0 = Inches(2.6)
gw  = Inches(3.85); gh = Inches(2.05)
gap = Inches(0.1)
for i, (emoji, title, body) in enumerate(features):
    col = i % 3; row = i // 3
    x = gx0 + col * (gw + gap)
    y = gy0 + row * (gh + gap)
    add_card(s, x, y, gw, gh, fill=WARM_WHITE, line=SAGE_LIGHT)
    add_text(s, emoji, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
             size=22)
    add_text(s, title, x + Inches(0.8), y + Inches(0.22), gw - Inches(1.0), Inches(0.4),
             font=SERIF, size=17, color=MOSS, bold=True)
    add_text(s, body, x + Inches(0.25), y + Inches(0.85), gw - Inches(0.5), gh - Inches(1.0),
             size=11, color=BARK, line_spacing=1.3)

add_page_number(s, 5, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Why now
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Why now",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Three forces have converged in the last 36 months.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

drivers = [
    ("Policy", "Denmark mandated 11 weeks earmarked paternity leave in 2022. Every Nordic country followed. Fathers are now structurally expected to be present."),
    ("Generation", "Millennials and Gen Z are the first generation of parents where 70%+ of fathers install the same pregnancy app their partner uses. They stop opening it because it ignores them."),
    ("Backlash", "Vulnerable-audience advertising and dark patterns face regulatory pressure. The window for a privacy-first, ad-free alternative is open."),
]
for i, (label, body) in enumerate(drivers):
    y = Inches(2.7 + i * 1.4)
    # Left label rail
    add_card(s, Inches(0.8), y, Inches(2.6), Inches(1.2), fill=MOSS)
    add_text(s, label, Inches(0.8), y, Inches(2.6), Inches(1.2),
             font=SERIF, size=22, color=TEXT_LIGHT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body card
    add_card(s, Inches(3.6), y, Inches(8.9), Inches(1.2), fill=WARM_WHITE, line=SAGE_LIGHT)
    add_text(s, body, Inches(3.85), y + Inches(0.18), Inches(8.4), Inches(0.95),
             size=12.5, color=BARK, line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

add_page_number(s, 6, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Market
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Market opportunity",
         Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Starting in Denmark, expanding into a €240M European SAM.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# 3 stacked cards: Denmark, Nordic, Europe
tiers = [
    ("Denmark",          "60K births / year",   "~€2–3M annual SAM",   "Year 0–1: launch market",     MOSS),
    ("Nordic",           "268K births / year",  "~€10–20M annual SAM", "Year 1–2: SE, NO, FI",        SAGE),
    ("Western Europe",   "3.5M births / year",  "~€120–240M annual SAM","Year 2+: DE, NL, UK, FR",   CLAY),
]
for i, (region, births, sam, when, color) in enumerate(tiers):
    y = Inches(2.7 + i * 1.45)
    add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.25), fill=WARM_WHITE, line=SAGE_LIGHT)
    # Color accent rail
    rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.18), Inches(1.25))
    rail.fill.solid(); rail.fill.fore_color.rgb = color
    rail.line.fill.background(); rail.shadow.inherit = False

    add_text(s, region, Inches(1.2), y + Inches(0.18), Inches(3.2), Inches(0.4),
             font=SERIF, size=20, color=MOSS, bold=True)
    add_text(s, when, Inches(1.2), y + Inches(0.65), Inches(3.2), Inches(0.4),
             size=11, color=STONE, italic=True)
    add_text(s, births, Inches(4.8), y + Inches(0.18), Inches(3.8), Inches(0.4),
             font=SERIF, size=18, color=BARK)
    add_text(s, "Births per year", Inches(4.8), y + Inches(0.7), Inches(3.8), Inches(0.3),
             size=9, color=STONE)
    add_text(s, sam, Inches(8.5), y + Inches(0.18), Inches(3.8), Inches(0.4),
             font=SERIF, size=18, color=CLAY, bold=True)
    add_text(s, "Annual addressable revenue", Inches(8.5), y + Inches(0.7), Inches(3.8), Inches(0.3),
             size=9, color=STONE)

add_page_number(s, 7, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Business model
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Business model",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Freemium family subscription. No ads. Never sold data.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# Price tiers
prices = [
    ("Free",      "Core app",        "Week-by-week · Sleep tracker · Diary · Partner sync · Basic chat", WARM_WHITE, MOSS),
    ("Monthly",   "39 DKK / mo",     "Premium for both parents — flexible, cancel anytime",              SAGE_LIGHT, MOSS),
    ("Annual",    "299 DKK / yr",    "≈25 DKK / mo — best value, anchors the choice",                    MOSS, TEXT_LIGHT),
    ("Bundle",    "499 DKK once",    "Full pregnancy + first 24 months. Lower friction.",                CLAY_LIGHT, BARK),
]
for i, (tier, price, body, fill, txt) in enumerate(prices):
    x = Inches(0.8 + i * 3.05)
    add_card(s, x, Inches(2.7), Inches(2.9), Inches(3.0), fill=fill,
             line=None if fill == MOSS else SAGE_LIGHT)
    add_text(s, tier, x + Inches(0.25), Inches(2.95), Inches(2.4), Inches(0.4),
             size=10, color=txt, bold=True)
    add_text(s, price, x + Inches(0.25), Inches(3.45), Inches(2.4), Inches(0.7),
             font=SERIF, size=26, color=txt, bold=True)
    add_text(s, body, x + Inches(0.25), Inches(4.4), Inches(2.4), Inches(1.3),
             size=11, color=txt, line_spacing=1.3)

# Adjacent revenue
add_text(s, "Year 2+ adjacent streams",
         Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
         font=SERIF, size=16, color=MOSS, bold=True)
add_text(s, "Affiliate commerce (curated wishlist) · Premium courses (99–299 DKK) · B2B / employer benefits · Municipal pilots",
         Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
         size=12, color=BARK)

add_page_number(s, 8, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Competition
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Competitive landscape",
         Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Big market, generic positioning. Two-parent design is our moat.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# Table header
headers = ["Player", "Region", "Model", "Gap we exploit"]
col_xs = [0.8, 4.0, 6.4, 8.8]
col_ws = [3.0, 2.2, 2.2, 3.7]
hy = Inches(2.65)
for i, h in enumerate(headers):
    add_text(s, h.upper(), Inches(col_xs[i]), hy, Inches(col_ws[i]), Inches(0.4),
             size=10, color=STONE, bold=True)
# divider
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.015))
line.fill.solid(); line.fill.fore_color.rgb = SAGE; line.line.fill.background()
line.shadow.inherit = False

rows = [
    ("Preglife",     "SE/DK/NO",  "Freemium ~39 SEK/mo",   "Mom-centric. Dads ignored."),
    ("Flo",          "Global",    "Freemium ~€7/mo",       "Cycle-tracking origin. Weak partner features."),
    ("BabyCenter",   "Global",    "Ad-supported",          "Heavy advertising. Trust issues."),
    ("Ovia",         "USA",       "B2B employer benefit",  "No EU consumer presence."),
    ("Peanut",       "Global",    "Freemium social",       "Community-driven. Not couple-focused."),
    ("Glow",         "Global",    "Freemium ~€7/mo",       "Cycle-first. No postpartum depth."),
]
for i, row in enumerate(rows):
    y = Inches(3.15 + i * 0.55)
    for j, val in enumerate(row):
        font_bold = (j == 0)
        color = MOSS if j == 0 else BARK
        size = 12.5 if j == 0 else 11.5
        add_text(s, val, Inches(col_xs[j]), y, Inches(col_ws[j]), Inches(0.45),
                 size=size, color=color, bold=font_bold,
                 anchor=MSO_ANCHOR.MIDDLE,
                 font=SERIF if j == 0 else SANS)

# Bottom moat callout
add_card(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.65), fill=MOSS)
add_text(s, "Melo's moat: structurally hard to copy — built around fathers from day one, with content depth competitors don't have.",
         Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.65),
         size=12, color=TEXT_LIGHT, italic=True, anchor=MSO_ANCHOR.MIDDLE)

add_page_number(s, 9, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Traction
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Traction",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Built. Ready. Launch-blockers being closed this quarter.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# Big number callouts
callouts = [
    ("35+",   "pages / screens built",        MOSS),
    ("36",    "weeks of pregnancy content",   SAGE),
    ("5",     "age-stages of dad content",    CLAY),
    ("GDPR",  "data architecture verified",   MOSS),
]
for i, (num, label, color) in enumerate(callouts):
    x = Inches(0.8 + i * 3.05)
    add_card(s, x, Inches(2.7), Inches(2.9), Inches(2.2), fill=WARM_WHITE, line=SAGE_LIGHT)
    add_text(s, num, x, Inches(2.85), Inches(2.9), Inches(1.0),
             font=SERIF, size=48, color=color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, label, x, Inches(4.0), Inches(2.9), Inches(0.5),
             size=11, color=BARK, align=PP_ALIGN.CENTER)

# Status checklist
add_text(s, "What's done",
         Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.4),
         font=SERIF, size=16, color=MOSS, bold=True)
add_lines(s, [
    "Feature-complete v1 build",
    "Postpartum dad experience (industry-first depth)",
    "Privacy policy + GDPR-compliant Row Level Security",
    "Brand, design system, app icon",
], Inches(0.8), Inches(5.65), Inches(5.7), Inches(2.0),
   size=12, color=BARK, bullet=True, line_spacing=1.4)

add_text(s, "In motion",
         Inches(7.0), Inches(5.2), Inches(5.5), Inches(0.4),
         font=SERIF, size=16, color=CLAY, bold=True)
add_lines(s, [
    "Clinical sign-off (Danish midwife) — in review",
    "App Store submission — pending sign-off",
    "Launch partner conversations (Mødrehjælpen)",
    "Internal TestFlight with first cohort",
], Inches(7.0), Inches(5.65), Inches(5.5), Inches(2.0),
   size=12, color=BARK, bullet=True, line_spacing=1.4)

add_page_number(s, 10, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Roadmap
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Roadmap",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Twelve months from launch to Nordic scale.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# 4 quarter columns
quarters = [
    ("Q3 2026", "Launch DK",
     ["Soft launch, free tier",
      "Clinical endorsement live",
      "Target 5K installs",
      "1.5K WAU end of Q"]),
    ("Q4 2026", "Monetize",
     ["Premium subscription live",
      "Target 3% conversion",
      "First media coverage",
      "250 paying families"]),
    ("Q1 2027", "Nordic",
     ["Swedish localization",
      "Soft launch SE",
      "First B2B pilot",
      "Series A prep begins"]),
    ("Q2 2027", "Scale",
     ["Norwegian launch",
      "25K installs across Nordic",
      "Series A close target",
      "Hiring: content + iOS"]),
]
for i, (q, title, items) in enumerate(quarters):
    x = Inches(0.8 + i * 3.05)
    # Header band
    add_card(s, x, Inches(2.7), Inches(2.9), Inches(0.8), fill=MOSS)
    add_text(s, q, x, Inches(2.7), Inches(2.9), Inches(0.4),
             size=11, color=TEXT_LIGHT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, title, x, Inches(3.05), Inches(2.9), Inches(0.4),
             font=SERIF, size=18, color=TEXT_LIGHT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body
    add_card(s, x, Inches(3.55), Inches(2.9), Inches(3.0), fill=WARM_WHITE, line=SAGE_LIGHT)
    add_lines(s, items, x + Inches(0.25), Inches(3.7), Inches(2.65), Inches(2.7),
              size=11.5, color=BARK, bullet=True, line_spacing=1.4, gap_after=3)

add_page_number(s, 11, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Team
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_brand_mark(s)

add_text(s, "Team",
         Inches(0.8), Inches(1.1), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=MOSS, bold=True)
add_text(s, "Lean by design. Next hires unlocked by this round.",
         Inches(0.8), Inches(1.85), Inches(11), Inches(0.4),
         size=14, color=STONE, italic=True)

# Founder card (large)
add_card(s, Inches(0.8), Inches(2.7), Inches(5.8), Inches(4.2), fill=WARM_WHITE, line=SAGE_LIGHT)
# Avatar placeholder circle
add_accent_dot(s, Inches(1.1), Inches(3.0), size=Inches(0.9), color=CLAY_LIGHT)
add_text(s, "LB", Inches(1.1), Inches(3.0), Inches(0.9), Inches(0.9),
         font=SERIF, size=28, color=CLAY, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "Louise Bencard", Inches(2.25), Inches(3.0), Inches(4), Inches(0.5),
         font=SERIF, size=22, color=MOSS, bold=True)
add_text(s, "Founder", Inches(2.25), Inches(3.5), Inches(4), Inches(0.4),
         size=12, color=STONE, italic=True)
add_text(s, "[Background — to fill in: design / product / clinical / parental story]",
         Inches(1.1), Inches(4.2), Inches(5.4), Inches(2.5),
         size=12, color=BARK, italic=True, line_spacing=1.4)

# Next hires
add_card(s, Inches(6.9), Inches(2.7), Inches(5.6), Inches(4.2), fill=SAGE_LIGHT)
add_text(s, "Next hires (this round)",
         Inches(7.1), Inches(2.95), Inches(5.3), Inches(0.5),
         font=SERIF, size=18, color=MOSS, bold=True)
add_lines(s, [
    "Senior content lead — clinical / editorial",
    "iOS engineer, full-time",
    "Part-time designer",
    "Customer support",
], Inches(7.1), Inches(3.55), Inches(5.3), Inches(2.5),
   size=13, color=BARK, bullet=True, line_spacing=1.5, gap_after=6)

add_text(s, "Plus a clinical advisor signing all content.",
         Inches(7.1), Inches(6.3), Inches(5.3), Inches(0.4),
         size=11, color=CLAY, italic=True)

add_page_number(s, 12, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — The Ask
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, MOSS)

# Decorative dots
for (x, y) in [(11.4, 0.9), (12.0, 1.3), (12.6, 1.0)]:
    add_accent_dot(s, Inches(x), Inches(y), size=Inches(0.12), color=CLAY)

add_brand_mark(s, on_dark=True)

add_text(s, "The ask",
         Inches(0.8), Inches(1.5), Inches(8), Inches(0.7),
         font=SERIF, size=40, color=TEXT_LIGHT, bold=True)
add_text(s, "Twelve months runway to validated launch and Nordic scale.",
         Inches(0.8), Inches(2.25), Inches(11.5), Inches(0.4),
         size=14, color=TEXT_LIGHT, italic=True)

# Big ask number
add_text(s, "€ [AMOUNT]",
         Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.6),
         font=SERIF, size=88, color=TEXT_LIGHT, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

# Use of funds bar
uses = [
    ("60%", "Team",       "Content lead, iOS engineer, designer, support"),
    ("15%", "Content",    "Clinical sign-off, professional video/illustration assets"),
    ("15%", "Marketing",  "Paid pilot, PR, clinical endorsement campaign"),
    ("10%", "Operations", "Legal (GDPR, App Store), accounting, working capital"),
]
for i, (pct, label, body) in enumerate(uses):
    x = Inches(0.8 + i * 3.05)
    add_card(s, x, Inches(5.0), Inches(2.9), Inches(1.95), fill=MOSS_LIGHT, line=None)
    add_text(s, pct, x + Inches(0.25), Inches(5.2), Inches(2.4), Inches(0.6),
             font=SERIF, size=26, color=TEXT_LIGHT, bold=True)
    add_text(s, label, x + Inches(0.25), Inches(5.75), Inches(2.4), Inches(0.4),
             font=SERIF, size=14, color=TEXT_LIGHT, bold=True)
    add_text(s, body, x + Inches(0.25), Inches(6.15), Inches(2.5), Inches(0.8),
             size=10, color=TEXT_LIGHT, line_spacing=1.2)

add_text(s, "Looking for patient capital with consumer-health or family-tech experience.",
         Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3),
         size=11, color=TEXT_LIGHT, italic=True)

add_page_number(s, 13, 14)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Thank you
# ──────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_bg(s, MOSS)

for (x, y) in [(11.6, 6.4), (12.2, 6.7)]:
    add_accent_dot(s, Inches(x), Inches(y), size=Inches(0.12), color=CLAY)

# Hero logo top-left
s.shapes.add_picture(LOGO_HERO, Inches(0.8), Inches(0.8),
                     height=Inches(1.8), width=Inches(1.8))

add_text(s, "Thank you.",
         Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5),
         font=SERIF, size=80, color=TEXT_LIGHT, bold=True)

add_text(s, "Let's build the calm hand that walks alongside every new family.",
         Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.6),
         font=SERIF, size=22, color=TEXT_LIGHT, italic=True)

add_text(s, "Louise Bencard   ·   Founder",
         Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
         size=14, color=TEXT_LIGHT)
add_text(s, "hej@meloparents.com   ·   meloparents.com",
         Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
         size=12, color=TEXT_LIGHT, italic=True)


# ── Save ─────────────────────────────────────────────────────────────────────
output = "/sessions/gracious-lucid-hopper/mnt/outputs/melo-pitch/Melo-pitch-deck.pptx"
prs.save(output)
print(f"Saved → {output}")
